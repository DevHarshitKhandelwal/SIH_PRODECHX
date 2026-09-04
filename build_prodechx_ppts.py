import os
import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def generate_sih_official_ppt():
    template_path = "SIH2026-IDEA-Presentation-Format.pptx"
    output_path = "PRODECHX_SIH2026_Official_Submission.pptx"
    
    prs = pptx.Presentation(template_path)
    
    # Ultra High Contrast Colors
    NAVY = RGBColor(10, 25, 47)          # #0A192F - Deep Navy
    DARK_BLUE = RGBColor(30, 58, 138)    # #1E3A8A - Royal Dark Blue
    DARK_SLATE = RGBColor(15, 23, 42)     # #0F172A - Slate 900
    SLATE_TEXT = RGBColor(15, 23, 42)     # #0F172A - Maximum Contrast Dark Text
    ACCENT_GREEN = RGBColor(16, 185, 129) # #10B981 - Emerald Green
    ACCENT_GOLD = RGBColor(217, 119, 6)   # #D97706 - Amber Gold
    WHITE = RGBColor(255, 255, 255)
    
    img_hero = "assets/hero_infrastructure.jpg"
    img_arch = "assets/architecture_diagram.jpg"
    img_shap = "assets/shap_explainability.jpg"
    img_impact = "assets/impact_infographic.jpg"

    def apply_slide_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(241, 245, 249) # Clean Slate-100 backdrop
        bg.line.color.rgb = RGBColor(241, 245, 249)
        spTree = slide.shapes._spTree
        spTree.remove(bg._element)
        spTree.insert(2, bg._element)

    def remove_template_footer(slide):
        for shape in slide.shapes:
            if shape.has_text_frame:
                if "Footer" in shape.name or "@SIH" in shape.text_frame.text or "Template" in shape.text_frame.text:
                    shape.text_frame.clear()

    def format_title(slide, title_text):
        for shape in slide.shapes:
            if shape.name == "Title 1":
                shape.left = Inches(0.5)
                shape.top = Inches(0.12)
                shape.width = Inches(12.2)
                shape.height = Inches(0.95)
                tf = shape.text_frame
                tf.word_wrap = True
                tf.clear()
                p = tf.paragraphs[0]
                p.text = title_text.upper()
                p.font.size = Pt(28)
                p.font.bold = True
                p.font.color.rgb = NAVY

    def format_team_oval(slide):
        for shape in slide.shapes:
            if "Oval" in shape.name:
                shape.left = Inches(0.5)
                shape.top = Inches(0.15)
                shape.width = Inches(2.2)
                shape.height = Inches(0.7)
                tf = shape.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = "TEAM PRODECHX"
                p.font.size = Pt(12)
                p.font.bold = True
                p.font.color.rgb = WHITE
                shape.fill.solid()
                shape.fill.fore_color.rgb = DARK_BLUE

    def add_metric_badge(slide, left, top, width, height, text_line1, text_line2):
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        badge.fill.solid()
        badge.fill.fore_color.rgb = NAVY
        badge.line.color.rgb = ACCENT_GOLD
        badge.line.width = Pt(1.5)
        
        tf = badge.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        p1.text = text_line1
        p1.font.size = Pt(13.5)
        p1.font.bold = True
        p1.font.color.rgb = ACCENT_GOLD
        
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.text = text_line2
        p2.font.size = Pt(12.5)
        p2.font.color.rgb = WHITE

    # -------------------------------------------------------------
    # Slide 1: Title Page
    # -------------------------------------------------------------
    slide1 = prs.slides[0]
    apply_slide_bg(slide1)
    remove_template_footer(slide1)
    
    for shape in slide1.shapes:
        if shape.name == "Title 7":
            shape.left = Inches(0.5)
            shape.top = Inches(0.25)
            shape.width = Inches(12.2)
            shape.height = Inches(0.8)
            tf = shape.text_frame
            tf.clear()
            p0 = tf.paragraphs[0]
            p0.text = "SMART INDIA HACKATHON 2026 — TITLE PAGE"
            p0.font.size = Pt(28)
            p0.font.bold = True
            p0.font.color.rgb = ACCENT_GOLD
            
        elif shape.name == "Subtitle 3":
            shape.left = Inches(0.5)
            shape.top = Inches(1.05)
            shape.width = Inches(6.1)
            shape.height = Inches(0.85)
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = "PRODECHX — Infrastructure Risk Intelligence Platform"
            p.font.size = Pt(19)
            p.font.bold = True
            p.font.color.rgb = NAVY
            
        elif shape.name == "TextBox 9":
            shape.left = Inches(0.5)
            shape.top = Inches(1.9)
            shape.width = Inches(6.1)
            shape.height = Inches(5.2)
            tf = shape.text_frame
            tf.word_wrap = True
            tf.clear()
            
            details = [
                ("Problem Statement ID: ", "SIH26103"),
                ("Problem Statement Title: ", "Predictive Infrastructure Monitoring Platform"),
                ("Theme: ", "Smart Automation  |  Category: Software"),
                ("Organization: ", "Ministry of Statistics & Programme Implementation (MoSPI)"),
                ("Department: ", "Data Informatics & Innovation Division (DIID)"),
                ("Team ID: ", "SIH2026-TEAM-PRODECHX"),
                ("Team Name: ", "PRODECHX Innovators (Registered on Portal)")
            ]
            
            for idx, (label, val) in enumerate(details):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                r1 = p.add_run()
                r1.text = label
                r1.font.bold = True
                r1.font.size = Pt(14)
                r1.font.color.rgb = NAVY
                
                r2 = p.add_run()
                r2.text = val
                r2.font.size = Pt(14)
                r2.font.color.rgb = SLATE_TEXT

    for shape in list(slide1.shapes):
        if shape.name == "Picture 4":
            sp = shape._element
            sp.getparent().remove(sp)

    if os.path.exists(img_hero):
        slide1.shapes.add_picture(img_hero, Inches(6.8), Inches(1.9), width=Inches(5.9))
        add_metric_badge(slide1, 6.8, 5.4, 5.9, 1.4, "5,815 PAIMANA DATA UPDATES  |  2,231 PROJECTS", "₹29.87 Lakh Crore Sanctioned Budget Monitored")

    # -------------------------------------------------------------
    # Slide 2: IDEA TITLE & PROPOSED SOLUTION
    # -------------------------------------------------------------
    slide2 = prs.slides[1]
    apply_slide_bg(slide2)
    remove_template_footer(slide2)
    format_title(slide2, "IDEA TITLE: PRODECHX PREDICTIVE PLATFORM")
    format_team_oval(slide2)

    for shape in slide2.shapes:
        if shape.name == "TextBox 8":
            shape.left = Inches(0.5)
            shape.top = Inches(1.25)
            shape.width = Inches(6.1)
            shape.height = Inches(5.8)
            tf = shape.text_frame
            tf.word_wrap = True
            tf.clear()
            
            sections = [
                ("Proposed Solution (Idea Description):", [
                    "AI/ML platform converting static PAIMANA PDFs into 2-month advance predictive early-warnings across 2,231 central infrastructure projects (₹29.87L Cr budget)."
                ]),
                ("How it Addresses the Problem:", [
                    "XGBoost v2.0 classifier @ 0.45 threshold catches 450 high-risk projects with 82.4% recall, providing 60 days for proactive intervention."
                ]),
                ("Innovation and Uniqueness:", [
                    "100% Server-Side Temporal Leakage Isolation.",
                    "SHAP Explainability for non-blackbox root cause drivers.",
                    "Grounded RAG Assistant delivering zero-hallucination Q&A with direct page PDF citations ([PAIMANA April 2026, p.54])."
                ])
            ]
            
            for idx, (heading, bullets) in enumerate(sections):
                p_head = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p_head.text = heading
                p_head.font.bold = True
                p_head.font.size = Pt(15.5)
                p_head.font.color.rgb = NAVY
                
                for b in bullets:
                    p = tf.add_paragraph()
                    p.text = "• " + b
                    p.font.size = Pt(13.5)
                    p.font.color.rgb = SLATE_TEXT

    if os.path.exists(img_hero):
        slide2.shapes.add_picture(img_hero, Inches(6.8), Inches(1.35), width=Inches(5.9))
        add_metric_badge(slide2, 6.8, 4.95, 5.9, 1.4, "60-DAY EARLY WARNING  |  82.4% RECALL", "Zero Hallucination Grounded AI Assistant")

    # -------------------------------------------------------------
    # Slide 3: TECHNICAL APPROACH
    # -------------------------------------------------------------
    slide3 = prs.slides[2]
    apply_slide_bg(slide3)
    remove_template_footer(slide3)
    format_title(slide3, "TECHNICAL APPROACH & PIPELINE")
    format_team_oval(slide3)

    for shape in slide3.shapes:
        if shape.name == "TextBox 8":
            shape.left = Inches(0.5)
            shape.top = Inches(1.25)
            shape.width = Inches(6.1)
            shape.height = Inches(5.8)
            tf = shape.text_frame
            tf.word_wrap = True
            tf.clear()
            
            sections = [
                ("Technologies to be Used:", [
                    "Frontend: Next.js 14, React, TypeScript, Tailwind CSS, Recharts.",
                    "Backend & DB: Node.js Express API, Supabase PostgreSQL + pgvector.",
                    "ML Microservice: Python 3.11 FastAPI, XGBoost, Scikit-Learn, SHAP.",
                    "Ingestion: PyMuPDF, pdfplumber, SHA-256 duplicate checksum."
                ]),
                ("Methodology & Implementation Process:", [
                    "1. Drag-and-Drop Ingestion: PDF parsing & table extraction.",
                    "2. Leakage-Free Features: Expenditure & progress gap ratios.",
                    "3. Predictive ML & SHAP: 2-month early warning delay risk scoring.",
                    "4. Grounded RAG Assistant: Vector search with 100% citations."
                ])
            ]
            
            for idx, (heading, bullets) in enumerate(sections):
                p_head = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p_head.text = heading
                p_head.font.bold = True
                p_head.font.size = Pt(15.5)
                p_head.font.color.rgb = NAVY
                
                for b in bullets:
                    p = tf.add_paragraph()
                    p.text = "• " + b
                    p.font.size = Pt(13.5)
                    p.font.color.rgb = SLATE_TEXT

    if os.path.exists(img_arch):
        slide3.shapes.add_picture(img_arch, Inches(6.8), Inches(1.35), width=Inches(5.9))
        add_metric_badge(slide3, 6.8, 4.95, 5.9, 1.4, "NEXT.JS 14  |  FASTAPI ML  |  PGVECTOR", "100% Server-Side Temporal Feature Pipeline")

    # -------------------------------------------------------------
    # Slide 4: FEASIBILITY AND VIABILITY
    # -------------------------------------------------------------
    slide4 = prs.slides[3]
    apply_slide_bg(slide4)
    remove_template_footer(slide4)
    format_title(slide4, "FEASIBILITY AND RISK MITIGATION")
    format_team_oval(slide4)

    for shape in slide4.shapes:
        if shape.name == "TextBox 8":
            shape.left = Inches(0.5)
            shape.top = Inches(1.25)
            shape.width = Inches(6.1)
            shape.height = Inches(5.8)
            tf = shape.text_frame
            tf.word_wrap = True
            tf.clear()
            
            sections = [
                ("Analysis of Feasibility:", [
                    "Validated on 5,815 monthly observations across 2,231 projects (2,030 baseline cohort) with 100% working live software prototype."
                ]),
                ("Potential Challenges and Risks:", [
                    "1. Risk of temporal data leakage in predictive scoring.",
                    "2. Risk of AI hallucination in policy query responses.",
                    "3. Layout variations & noise in legacy PAIMANA PDF reports."
                ]),
                ("Strategies for Overcoming Challenges:", [
                    "Strict server-side feature construction isolating pre-period metrics.",
                    "Deterministic RAG grounding policy requiring 100% page citations.",
                    "Multi-stage table parser with PyMuPDF OCR fallback & review queue."
                ])
            ]
            
            for idx, (heading, bullets) in enumerate(sections):
                p_head = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p_head.text = heading
                p_head.font.bold = True
                p_head.font.size = Pt(15.5)
                p_head.font.color.rgb = NAVY
                
                for b in bullets:
                    p = tf.add_paragraph()
                    p.text = "• " + b
                    p.font.size = Pt(13.5)
                    p.font.color.rgb = SLATE_TEXT

    if os.path.exists(img_shap):
        slide4.shapes.add_picture(img_shap, Inches(6.8), Inches(1.35), width=Inches(5.9))
        add_metric_badge(slide4, 6.8, 4.95, 5.9, 1.4, "SHAP AUDIT  |  NON-BLACKBOX EXPLAINABILITY", "Transparent Risk Gauges & Root-Cause Attributions")

    # -------------------------------------------------------------
    # Slide 5: IMPACT AND BENEFITS
    # -------------------------------------------------------------
    slide5 = prs.slides[4]
    apply_slide_bg(slide5)
    remove_template_footer(slide5)
    format_title(slide5, "IMPACT AND BENEFITS")
    format_team_oval(slide5)

    for shape in slide5.shapes:
        if shape.name == "TextBox 8":
            shape.left = Inches(0.5)
            shape.top = Inches(1.25)
            shape.width = Inches(6.1)
            shape.height = Inches(5.8)
            tf = shape.text_frame
            tf.word_wrap = True
            tf.clear()
            
            sections = [
                ("Impact on Target Audience (MoSPI & Ministries):", [
                    "60-Day Lead Time: Gives decision-makers 2 months to resolve bottlenecks before major cost escalation occurs.",
                    "Automated Workflows: Replaces manual review of 100+ page PDFs with automated risk gauges and SHAP driver cards."
                ]),
                ("Benefits of Solution (Social, Economic, Policy):", [
                    "Economic Benefits: Saves thousands of crores in public infrastructure cost overruns.",
                    "Policy Benefits: Standardizes progress scoring across Railways, Highways, Power, Petroleum, and Telecom.",
                    "Social Benefits: Accelerates critical public utility delivery (highways, clean energy, rail corridors)."
                ])
            ]
            
            for idx, (heading, bullets) in enumerate(sections):
                p_head = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p_head.text = heading
                p_head.font.bold = True
                p_head.font.size = Pt(15.5)
                p_head.font.color.rgb = NAVY
                
                for b in bullets:
                    p = tf.add_paragraph()
                    p.text = "• " + b
                    p.font.size = Pt(13.5)
                    p.font.color.rgb = SLATE_TEXT

    if os.path.exists(img_impact):
        slide5.shapes.add_picture(img_impact, Inches(6.8), Inches(1.35), width=Inches(5.9))
        add_metric_badge(slide5, 6.8, 4.95, 5.9, 1.4, "NATIONAL POLICY IMPACT  |  FISCAL SAVINGS", "Standardized Cross-Ministry Performance Scoring")

    # -------------------------------------------------------------
    # Slide 6: RESEARCH AND REFERENCES
    # -------------------------------------------------------------
    slide6 = prs.slides[5]
    apply_slide_bg(slide6)
    remove_template_footer(slide6)
    format_title(slide6, "RESEARCH AND REFERENCES")
    format_team_oval(slide6)

    for shape in slide6.shapes:
        if shape.name == "TextBox 8":
            shape.left = Inches(0.5)
            shape.top = Inches(1.25)
            shape.width = Inches(12.2)
            shape.height = Inches(5.8)
            tf = shape.text_frame
            tf.word_wrap = True
            tf.clear()
            
            sections = [
                ("Authoritative Data Corpus & Research Work:", [
                    "MoSPI PAIMANA Flash Reports (April, May, June 2026): 5,815 monthly observations across 2,231 central sector infrastructure projects totaling ₹29.87 Lakh Cr sanctioned cost."
                ]),
                ("Empirical Model Validation Results:", [
                    "82.4% Recall @ 0.45 Operating Threshold (Random Forest / XGBoost v2.0 classifier).",
                    "100% Citation Grounding Accuracy & 0% Hallucination Rate on unsupported user queries."
                ]),
                ("Academic & Industry Technical Standards / Links:", [
                    "SHapley Additive exPlanations (SHAP) — Lundberg & Lee (2017) for non-blackbox feature attributions.",
                    "PostgreSQL pgvector extension for dense embedding storage & HNSW similarity search.",
                    "MoSPI Infrastructure & Project Monitoring Division (IPMD) Guidelines & SIH26103 Problem Statement."
                ])
            ]
            
            for idx, (heading, bullets) in enumerate(sections):
                p_head = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p_head.text = heading
                p_head.font.bold = True
                p_head.font.size = Pt(16)
                p_head.font.color.rgb = NAVY
                
                for b in bullets:
                    p = tf.add_paragraph()
                    p.text = "• " + b
                    p.font.size = Pt(14)
                    p.font.color.rgb = SLATE_TEXT

    # Remove Slide 7 if present
    if len(prs.slides) >= 7:
        rId = prs.slides._sldIdLst[6].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[6]

    prs.save(output_path)
    print(f"Saved template footer-free, large font official SIH submission presentation to {output_path}")

def generate_detailed_proposal_ppt():
    output_path = "PRODECHX_Detailed_Project_Proposal.pptx"
    
    prs = pptx.Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    NAVY = RGBColor(10, 25, 47)
    DARK_BLUE = RGBColor(30, 58, 138)
    SLATE = RGBColor(15, 23, 42)
    GRAY_TEXT = RGBColor(71, 85, 105)
    LIGHT_BG = RGBColor(248, 250, 252)
    WHITE = RGBColor(255, 255, 255)
    BORDER_COLOR = RGBColor(203, 213, 225)
    ACCENT_GREEN = RGBColor(16, 185, 129)
    ACCENT_GOLD = RGBColor(217, 119, 6)
    
    img_hero = "assets/hero_infrastructure.jpg"
    img_arch = "assets/architecture_diagram.jpg"
    img_shap = "assets/shap_explainability.jpg"
    img_impact = "assets/impact_infographic.jpg"

    def add_slide_header(slide, title_text, category_text="PRODECHX PROJECT PROPOSAL"):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(241, 245, 249)
        bg.line.color.rgb = RGBColor(241, 245, 249)
        
        banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.15))
        banner.fill.solid()
        banner.fill.fore_color.rgb = NAVY
        banner.line.color.rgb = NAVY
        
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.15), Inches(13.333), Inches(0.06))
        strip.fill.solid()
        strip.fill.fore_color.rgb = ACCENT_GOLD
        strip.line.color.rgb = ACCENT_GOLD
        
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.12), Inches(11.5), Inches(0.3))
        tf_cat = cat_box.text_frame
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(12)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_GREEN
        
        t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.38), Inches(11.5), Inches(0.6))
        tf_t = t_box.text_frame
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(25)
        p_t.font.bold = True
        p_t.font.color.rgb = WHITE
        
        footer = slide.shapes.add_textbox(Inches(0.8), Inches(7.08), Inches(11.733), Inches(0.35))
        tf_f = footer.text_frame
        p_f = tf_f.paragraphs[0]
        p_f.text = "PRODECHX | MoSPI Predictive Infrastructure Monitoring Platform (SIH26103)"
        p_f.font.size = Pt(11.5)
        p_f.font.color.rgb = GRAY_TEXT

    def add_card(slide, left, top, width, height, title, content_bullets, bg_color=LIGHT_BG, border_color=BORDER_COLOR, title_color=NAVY):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.2)
        
        tb = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.15), Inches(width - 0.4), Inches(height - 0.3))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = title_color
        
        for bullet in content_bullets:
            p_b = tf.add_paragraph()
            p_b.text = "• " + bullet
            p_b.font.size = Pt(13.5)
            p_b.font.color.rgb = SLATE

    # Slide 1: Title Slide
    s1 = prs.slides.add_slide(blank_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = NAVY
    bg1.line.color.rgb = NAVY
    
    tb1 = s1.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.2), Inches(4.8))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "SMART INDIA HACKATHON 2026 PROPOSAL"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD
    
    p2 = tf1.add_paragraph()
    p2.text = "PRODECHX"
    p2.font.size = Pt(52)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    
    p3 = tf1.add_paragraph()
    p3.text = "Predictive Infrastructure Monitoring, Risk & Early-Warning Intelligence Platform"
    p3.font.size = Pt(22)
    p3.font.color.rgb = RGBColor(191, 219, 254)
    
    p4 = tf1.add_paragraph()
    p4.text = "\nMinistry of Statistics and Programme Implementation (MoSPI)\nData Informatics & Innovation Division (DIID)\nProblem Statement ID: SIH26103  |  Category: Software"
    p4.font.size = Pt(14.5)
    p4.font.color.rgb = WHITE
    
    if os.path.exists(img_hero):
        s1.shapes.add_picture(img_hero, Inches(7.2), Inches(1.8), width=Inches(5.4))

    # Slide 2: Executive Overview
    s2 = prs.slides.add_slide(blank_layout)
    add_slide_header(s2, "1. Executive Summary & Core Vision")
    add_card(s2, 0.8, 1.45, 5.7, 5.3, "Executive Summary", [
        "PRODECHX is an AI/ML-powered infrastructure monitoring platform for MoSPI (DIID).",
        "Transforms static monthly PAIMANA PDF reporting into a 2-month advance predictive early-warning system.",
        "Monitors ₹29.87 Lakh Cr central infrastructure portfolio across 2,231 projects in India.",
        "Answers critical policy questions: Which projects face cost overruns? Why are they delayed? What interventions are needed?"
    ])
    if os.path.exists(img_hero):
        s2.shapes.add_picture(img_hero, Inches(6.8), Inches(1.8), width=Inches(5.7))

    # Slide 3: Need of Project
    s3 = prs.slides.add_slide(blank_layout)
    add_slide_header(s3, "2. Need of Project")
    add_card(s3, 0.8, 1.45, 3.7, 5.3, "Current Challenges", [
        "Reactive Risk Detection: Overruns detected AFTER budget is spent.",
        "Unstructured PDF Data: 100+ page PAIMANA PDFs slow down analysis.",
        "Information Silos: Lack of cross-ministry performance benchmarks."
    ], bg_color=RGBColor(254, 242, 242), border_color=RGBColor(252, 165, 165), title_color=RGBColor(185, 28, 28))
    
    add_card(s3, 4.8, 1.45, 3.7, 5.3, "The PRODECHX Solution", [
        "60-Day Lead Time: Early warning flags enable proactive intervention.",
        "Automated Ingestion: Intelligent PDF parsing extracts tables & text.",
        "SHAP Explainability: Clear root-cause attributions for delay drivers."
    ], bg_color=RGBColor(240, 253, 244), border_color=RGBColor(134, 239, 172), title_color=RGBColor(21, 128, 61))
    
    add_card(s3, 8.8, 1.45, 3.7, 5.3, "Strategic MoSPI Impact", [
        "Fiscal Discipline: Prevents thousand-crore budget overruns.",
        "Evidence-Based Policy: Grounded citations link scores to PDF source.",
        "National Visibility: Executive dashboards for senior leadership."
    ], bg_color=RGBColor(239, 246, 255), border_color=RGBColor(147, 197, 253), title_color=NAVY)

    # Slide 4: Functionality Overview
    s4 = prs.slides.add_slide(blank_layout)
    add_slide_header(s4, "3. Functionality — System User Personas")
    add_card(s4, 0.8, 1.45, 5.7, 2.5, "Section A: National Administrator", [
        "National portfolio overview & risk heatmap",
        "Cross-ministry performance benchmarking",
        "Systemic delay driver analytics & early warnings"
    ])
    add_card(s4, 6.8, 1.45, 5.7, 2.5, "Section B: Ministry Officer", [
        "Ministry-specific project filtering & monitoring",
        "Cost escalation & schedule delay alerts",
        "Intervention action assignment & status tracking"
    ])
    add_card(s4, 0.8, 4.15, 5.7, 2.5, "Section C: Project Officer (Nodal)", [
        "Monthly project update & milestone logging",
        "PAIMANA PDF report uploads & validation",
        "Field progress commentary & bottleneck logging"
    ])
    add_card(s4, 6.8, 4.15, 5.7, 2.5, "Section D: Data Analyst & ML Engineer", [
        "Model accuracy, confusion matrix & threshold tuning",
        "Global SHAP feature attributions & leakage audits",
        "Data export (CSV/Excel) & ingestion audit logs"
    ])

    # Slide 5: Functionality - Admin
    s5 = prs.slides.add_slide(blank_layout)
    add_slide_header(s5, "4. Functionality — Section A: National Admin")
    add_card(s5, 0.8, 1.45, 11.733, 5.3, "National Administrator Capability Matrix", [
        "Portfolio Risk Matrix: Executive view of total budget (₹29.87L Cr) and calibrated high-risk project cohort (450 projects / 22.2%).",
        "Sector & Ministry Analytics: Comparative breakdown across Road Transport, Railways, Power, Petroleum, and Telecom.",
        "Early-Warning Command Center: Instant notification feed for projects crossing critical risk thresholds (>0.45 risk score).",
        "Data Quality & System Health: Monitoring PDF ingestion status, checksum verifications, and manual review queues.",
        "Audit Trail & Compliance: Immutable logging of user activities, score re-calibrations, and report exports."
    ])

    # Slide 6: Functionality - Ministry Officer
    s6 = prs.slides.add_slide(blank_layout)
    add_slide_header(s6, "5. Functionality — Section B: Ministry Officer")
    add_card(s6, 0.8, 1.45, 11.733, 5.3, "Ministry Officer Capability Matrix", [
        "Ministry Project Register: Filtered view of ongoing, delayed, and high-risk projects under the specific ministry.",
        "Milestone Tracking & Gap Analysis: Monitoring physical progress vs financial expenditure ratios to spot red flags early.",
        "Cost Overrun Deep-Dive: Detailed view of original vs revised cost estimates and target completion dates.",
        "Intervention Action Allocation: Assigning mitigation tasks to field project officers with clear deadlines.",
        "Custom PDF Report Generation: Exporting ministry risk summaries for cabinet and committee reviews."
    ])

    # Slide 7: Functionality - Project Officer
    s7 = prs.slides.add_slide(blank_layout)
    add_slide_header(s7, "6. Functionality — Section C: Project Officer")
    add_card(s7, 0.8, 1.45, 11.733, 5.3, "Project Officer Capability Matrix", [
        "Monthly Progress Submission: Direct input interface for physical progress %, expenditure, and completion dates.",
        "PAIMANA Document Submission: Drag-and-drop upload of official monthly PAIMANA Flash Report PDFs.",
        "Automated Validation Check: Ensures total expenditure does not exceed approved estimates without flag.",
        "Field Bottleneck Logging: Documenting land acquisition, forest clearance, or contractor issues causing delay.",
        "Mitigation Status Updates: Updating progress on administrative interventions assigned by ministry leadership."
    ])

    # Slide 8: Functionality - Data Analyst
    s8 = prs.slides.add_slide(blank_layout)
    add_slide_header(s8, "7. Functionality — Section D: Data Analyst")
    add_card(s8, 0.8, 1.45, 11.733, 5.3, "Data Analyst & ML Capability Matrix", [
        "ML Model Diagnostics: Inspection of Random Forest / XGBoost model metrics (82.4% Recall, Precision, ROC-AUC curve).",
        "SHAP Feature Attribution Explorer: Visualizing global feature rankings (expenditure ratio, progress gap, time elapsed).",
        "Threshold Tuning Workbench: Testing model operating thresholds (evaluating impact of threshold adjustments).",
        "Data Ingestion & Integrity Audits: Reviewing extracted tables, duplicate detection logs, and review queue items.",
        "Export Engine: Exporting complete cleansed datasets, feature dictionaries, and prediction results in CSV format."
    ])

    # Slide 9: Modules - Overview & Dashboard
    s9 = prs.slides.add_slide(blank_layout)
    add_slide_header(s9, "8. Modules / Web Pages — Executive Dashboard")
    add_card(s9, 0.8, 1.45, 5.7, 5.3, "Module 1: Executive Overview Dashboard", [
        "Route: / (Overview Page)",
        "Key Metrics: Total Sanctioned Cost (₹29.87L Cr), Total Projects (2,231), High Risk Cohort (450), Baseline Cohort (2,030).",
        "Visual Widgets: Risk score distribution chart, sector risk bar chart, top delayed projects carousel.",
        "Action Filters: Instant toggle between All Projects, High Risk, Medium Risk, and Low Risk."
    ])
    if os.path.exists(img_hero):
        s9.shapes.add_picture(img_hero, Inches(6.8), Inches(1.8), width=Inches(5.7))

    # Slide 10: Modules - Projects Register
    s10 = prs.slides.add_slide(blank_layout)
    add_slide_header(s10, "9. Modules / Web Pages — Projects Register")
    add_card(s10, 0.8, 1.45, 11.733, 5.3, "Module 3: Master Projects Register & Smart Search", [
        "Route: /projects (Projects Master Page)",
        "High-Density Grid: TanStack Table displaying Project Code, Name, Ministry, Sector, Cost, Physical Progress, and Risk Level.",
        "Multi-Faceted Search & Filter: Client-side search by project ID (e.g. '612786'), filter by Ministry or Risk Level.",
        "Baseline Cohort Indicator: Clear labeling distinguishing eligible T1 baseline projects from newly added T0 projects.",
        "Export Capabilities: One-click export of filtered tabular data to CSV format for offline reporting."
    ])

    # Slide 11: Modules - Deep Project Audit
    s11 = prs.slides.add_slide(blank_layout)
    add_slide_header(s11, "10. Modules / Web Pages — Deep Project Audit")
    add_card(s11, 0.8, 1.45, 5.7, 5.3, "Module 4: Deep Project Audit & SHAP Explainability", [
        "Route: /projects/[id] (Project Audit Page)",
        "API Security Architecture: Client calls POST /predict/project with project_id; server constructs features 100% server-side.",
        "Risk Gauge & Probabilities: Predicted Risk Score (84/100) and Cost Overrun Probability (78.2%).",
        "SHAP Waterfall Attributions: Transparent breakdown of top risk drivers.",
        "Linked Document Citation: Direct button linking to exact PAIMANA PDF source report and page number."
    ])
    if os.path.exists(img_shap):
        s11.shapes.add_picture(img_shap, Inches(6.8), Inches(1.8), width=Inches(5.7))

    # Slide 12: Modules - Document Management
    s12 = prs.slides.add_slide(blank_layout)
    add_slide_header(s12, "11. Modules / Web Pages — Document Management")
    add_card(s12, 0.8, 1.45, 11.733, 5.3, "Module 5: PAIMANA Document Management & Ingestion", [
        "Route: /admin/documents (Document Management Page)",
        "Drag-and-Drop Uploader: Supports bulk PAIMANA PDF uploads with reporting month/year tagging.",
        "Processing Pipeline Monitor: Real-time status tracking (Uploaded -> Extracting -> Validating -> Completed).",
        "Duplicate Detection Engine: Computes SHA-256 checksum to prevent duplicate report processing.",
        "Manual Review Workbench: Interface for resolving ambiguous project title matches or missing metadata fields."
    ])

    # Slide 13: Modules - RAG AI Assistant
    s13 = prs.slides.add_slide(blank_layout)
    add_slide_header(s13, "12. Modules / Web Pages — Grounded RAG Assistant")
    add_card(s13, 0.8, 1.45, 11.733, 5.3, "Module 6: Grounded RAG AI Assistant Workspace", [
        "Route: /assistant (AI Assistant Workspace Page)",
        "Dual-Panel Layout: Left panel for conversation feed; Right panel for Cited Evidence Vault.",
        "Grounded Querying: Instant natural language search over 5,815 project records and PAIMANA PDFs.",
        "Strict Citation Policy: Every response includes explicit source tags (e.g. [PAIMANA April 2026, p.54]).",
        "Zero-Hallucination Policy: Explicitly responds with 'insufficient evidence' for out-of-scope prompts."
    ])

    # Slide 14: Database Structure Overview
    s14 = prs.slides.add_slide(blank_layout)
    add_slide_header(s14, "13. Database Structure — Architecture & ER Overview")
    add_card(s14, 0.8, 1.45, 5.7, 5.3, "Relational & Vector Architecture", [
        "Engine: PostgreSQL (Supabase Cloud).",
        "Vector Extension: pgvector for dense embedding storage and HNSW vector similarity indexes.",
        "Temporal History Model: Single project entity linked to multiple monthly project_updates records.",
        "Data Provenance: Every document_chunk links back to source document ID and page number."
    ])
    if os.path.exists(img_arch):
        s14.shapes.add_picture(img_arch, Inches(6.8), Inches(1.8), width=Inches(5.7))

    # Slide 15: Database Schemas Table
    s15 = prs.slides.add_slide(blank_layout)
    add_slide_header(s15, "14. Database Structure — Schema Specification")
    
    rows, cols = 6, 5
    left, top, width, height = Inches(0.8), Inches(1.45), Inches(11.733), Inches(5.2)
    table_shape = s15.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    headers = ["Table Name", "Primary Key", "Foreign Keys", "Key Fields", "Description"]
    col_widths = [Inches(2.2), Inches(1.5), Inches(2.0), Inches(3.2), Inches(2.833)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w
        
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13.5)
            p.font.bold = True
            p.font.color.rgb = WHITE
            
    table_data = [
        ("tbl_projects", "project_id", "ministry_id, sector_id", "project_code, original_cost, target_date", "Master project record catalogue"),
        ("tbl_project_updates", "update_id", "project_id, document_id", "report_month, physical_pct, expenditure", "Monthly PAIMANA status updates"),
        ("tbl_documents", "document_id", "uploaded_by", "file_name, storage_path, checksum, status", "Ingested PDF document metadata"),
        ("tbl_document_chunks", "chunk_id", "document_id", "page_number, content, embedding", "Vector embeddings for RAG QA"),
        ("tbl_ml_predictions", "prediction_id", "project_id", "risk_score, cost_prob, shap_values", "Stored prediction outcomes & SHAP")
    ]
    
    for i, row in enumerate(table_data, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG if i % 2 == 1 else WHITE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12.5)
                p.font.color.rgb = SLATE

    # Slide 16: Machine Learning Architecture
    s16 = prs.slides.add_slide(blank_layout)
    add_slide_header(s16, "15. Machine Learning Architecture & SHAP")
    add_card(s16, 0.8, 1.45, 5.7, 5.3, "Predictive Model & SHAP Specification", [
        "Model Type: Calibrated Random Forest / XGBoost v2.0 classifier.",
        "Prediction Horizon: 2-month advance early warning.",
        "Operating Threshold: 0.45 probability threshold (82.4% Recall).",
        "Explainability: SHAP local and global feature importance attributions.",
        "Leakage Prevention: Server-side feature engineering strictly isolating pre-period historical metrics."
    ])
    if os.path.exists(img_shap):
        s16.shapes.add_picture(img_shap, Inches(6.8), Inches(1.8), width=Inches(5.7))

    # Slide 17: Reports & Impact
    s17 = prs.slides.add_slide(blank_layout)
    add_slide_header(s17, "16. Reports, Decision Support & Key Metrics")
    add_card(s17, 0.8, 1.45, 5.7, 5.3, "Automated Reports & Value Proposition", [
        "Executive Risk Digest: High-level PDF/Excel summaries for MoSPI decision makers.",
        "Ministry Scorecards: Monthly risk ranking and progress compliance reports.",
        "Intervention Trackers: Open vs resolved action items across nodal project officers.",
        "PDF Lineage Reports: Audit trail mapping every data point back to source PAIMANA PDF page."
    ])
    if os.path.exists(img_impact):
        s17.shapes.add_picture(img_impact, Inches(6.8), Inches(1.8), width=Inches(5.7))

    # Slide 18: Tech Stack & System Specs
    s18 = prs.slides.add_slide(blank_layout)
    add_slide_header(s18, "17. Technology Stack & Specifications")
    add_card(s18, 0.8, 1.45, 5.7, 5.3, "Software Technology Stack", [
        "Frontend: Next.js 14, React, TypeScript, Tailwind CSS, Recharts.",
        "Main Backend: Node.js, Express.js, Supabase REST API, JWT RBAC.",
        "ML Service: Python 3.11, FastAPI, XGBoost, Scikit-learn, SHAP, Pandas.",
        "Database & Vector: Supabase PostgreSQL + pgvector extension.",
        "PDF Ingestion: PyMuPDF (fitz), pdfplumber, Tesseract OCR fallback.",
        "Deployment: Docker containers, Docker Compose, Cloud / Vercel."
    ])
    add_card(s18, 6.8, 1.45, 5.7, 5.3, "System Requirements & Performance", [
        "Hardware Specs: Server with 8 vCPUs, 16GB RAM, 100GB SSD.",
        "Client Specs: Any modern web browser (Chrome, Edge, Firefox, Safari).",
        "Ingestion Capacity: Up to 500 PAIMANA PDF reports processed per hour.",
        "API Throughput: Sub-100ms response time for prediction endpoints and vector search.",
        "SIH Winning Factor: 100% working prototype with empirical MoSPI data."
    ])

    prs.save(output_path)
    print(f"Saved template footer-free, large font detailed project proposal presentation to {output_path}")

if __name__ == "__main__":
    generate_sih_official_ppt()
    generate_detailed_proposal_ppt()
