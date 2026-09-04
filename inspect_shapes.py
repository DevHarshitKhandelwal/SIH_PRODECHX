import pptx

prs = pptx.Presentation('SIH2026-IDEA-Presentation-Format.pptx')
print("Slide width:", prs.slide_width.inches, "height:", prs.slide_height.inches)

for idx, slide in enumerate(list(prs.slides)[:6]):
    print(f"\n=== Slide {idx+1} ===")
    for s in slide.shapes:
        left = round(s.left.inches, 2) if hasattr(s, 'left') and s.left is not None else 0
        top = round(s.top.inches, 2) if hasattr(s, 'top') and s.top is not None else 0
        width = round(s.width.inches, 2) if hasattr(s, 'width') and s.width is not None else 0
        height = round(s.height.inches, 2) if hasattr(s, 'height') and s.height is not None else 0
        text = s.text_frame.text.replace('\n', ' ')[:50] if s.has_text_frame else ''
        print(f"  [{s.name}] id:{s.shape_id} pos=({left}, {top}) size=({width}x{height}) text='{text}'")
